#!/usr/bin/env python3
"""
Run ONCE to build the FAISS vector store from SCD training materials.
Re-run only when source documents change.

Usage:
    python scripts/ingest.py
"""
import json
import os
import re
import sys

import faiss
import numpy as np
from dotenv import load_dotenv

load_dotenv()
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from app.config import (
    TRAINING_MATERIALS_DIR,
    FAISS_INDEX_PATH,
    OPENAI_API_KEY,
    OCR_ENABLED,
    OCR_LANG,
    OCR_DPI,
    TESSERACT_CMD,
)
from openai import OpenAI

EMBED_MODEL = "text-embedding-3-small"
CHUNK_SIZE = 400
CHUNK_OVERLAP = 40
EMBED_BATCH_SIZE = 100

# OCR only contributes a line if the text layer is missing it. A page whose text
# layer already has at least this many characters still gets OCR'd, but only the
# net-new lines (e.g. text inside diagrams) are merged in.
_OCR_MIN_LINE_LEN = 4        # ignore tiny OCR fragments / noise
_OCR_DEDUP_LINE_LEN = 8      # lines >= this are checked against the text layer to avoid re-adding prose

client = OpenAI(api_key=OPENAI_API_KEY)

_ocr_ready = False


def _init_ocr():
    """Return the pytesseract module if OCR is usable, else None."""
    global _ocr_ready
    if not OCR_ENABLED:
        return None
    try:
        import pytesseract
        if TESSERACT_CMD:
            pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
        # Probe the binary so we fail loudly here, not mid-document.
        pytesseract.get_tesseract_version()
        import fitz  # noqa: F401  (PyMuPDF, for rendering)
        from PIL import Image  # noqa: F401
        if not _ocr_ready:
            print(f"  OCR enabled (lang={OCR_LANG}, dpi={OCR_DPI})")
            _ocr_ready = True
        return pytesseract
    except Exception as e:
        print(f"  WARNING: OCR requested but unavailable ({e}). Falling back to text layer only.")
        return None


def _normalize(s: str) -> str:
    return " ".join("".join(c.lower() if c.isalnum() else " " for c in s).split())


def _merge_ocr(text_layer: str, ocr_text: str) -> str:
    """Keep the clean text layer; append only OCR lines it does not already contain."""
    base = (text_layer or "").strip()
    base_norm = _normalize(base)
    extra: list[str] = []
    for raw in (ocr_text or "").splitlines():
        line = raw.strip()
        if len(line) < _OCR_MIN_LINE_LEN:
            continue
        norm = _normalize(line)
        if not norm:
            continue
        # Skip lines already present in the text layer (avoids duplicating prose).
        if len(line) >= _OCR_DEDUP_LINE_LEN and norm in base_norm:
            continue
        extra.append(line)
    if not extra:
        return base
    return (base + "\n" + "\n".join(extra)).strip() if base else "\n".join(extra)


def _ocr_page(pytesseract, page) -> str:
    import io
    from PIL import Image
    pix = page.get_pixmap(dpi=OCR_DPI)
    img = Image.open(io.BytesIO(pix.tobytes("png")))
    return pytesseract.image_to_string(img, lang=OCR_LANG)


def extract_pdf(path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)
    text_pages = [page.extract_text() or "" for page in reader.pages]

    pyt = _init_ocr()
    if pyt is None:
        return "\n".join(text_pages)

    import fitz
    doc = fitz.open(path)
    merged_pages: list[str] = []
    added = 0
    for i, text_layer in enumerate(text_pages):
        ocr_text = ""
        if i < doc.page_count:
            try:
                ocr_text = _ocr_page(pyt, doc[i])
            except Exception as e:
                print(f"    OCR failed on page {i + 1}: {e}")
        before = len(text_layer)
        merged = _merge_ocr(text_layer, ocr_text)
        if len(merged) > before:
            added += 1
        merged_pages.append(merged)
    doc.close()
    print(f"    OCR added content on {added}/{len(text_pages)} page(s)")
    return "\n".join(merged_pages)


def extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


# --- Cadre tagging -------------------------------------------------------------
# Canonical cadre codes: ASHA, ANM, MPW, CHO, SN, LT, MO, COUNSELLOR, plus "all".
# File-level mapping for the manuals; the roles PPTX is tagged per slide below.
SOURCE_CADRES = {
    "GoI_CHO SCD Training Manual (English).pdf": ["CHO"],
    "GoI_MO SCD Training Manual (English).pdf": ["MO"],
    "GoI_MPW and Asha SCD Training Manual (English).pdf": ["MPW", "ASHA"],
    "GoI_Staff Nurse SCD Training Manual (English).pdf": ["SN"],
    "GoI_Counseling Module (English).pdf": ["COUNSELLOR", "all"],
    "GoI_NSCAEM Operational Guidelines.pdf": ["all"],
}

_CADRE_KEYWORDS = [
    ("ASHA", ["accredited social health", "asha"]),
    ("ANM", ["auxiliary nurse", "anm"]),
    ("MPW", ["multi-purpose", "multipurpose", "mpw"]),
    ("CHO", ["community health officer", "cho"]),
    ("SN", ["staff nurse"]),
    ("LT", ["laboratory technician", "lab technician"]),
    ("MO", ["medical officer"]),
    ("COUNSELLOR", ["counsellor", "counselor"]),
]


def _detect_cadres(text: str) -> list[str]:
    t = (text or "").lower()
    hits = [code for code, kws in _CADRE_KEYWORDS if any(k in t for k in kws)]
    # A slide naming many cadres (e.g. the cross-cadre reference chart) is general.
    if len(hits) >= 3 or not hits:
        return ["all"]
    return hits


_DIAGRAM_DATA_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData"
)


def _slide_smartart_text(slide) -> str:
    """Extract SmartArt text from a slide — python-pptx can't read diagrams as
    normal shapes; the text lives in the related diagram data part (dataN.xml)."""
    runs: list[str] = []
    for rel in slide.part.rels.values():
        if rel.reltype == _DIAGRAM_DATA_RELTYPE:
            xml = rel.target_part.blob.decode("utf-8", "ignore")
            for t in re.findall(r"<a:t>(.*?)</a:t>", xml, re.S):
                s = t.strip()
                if s:
                    runs.append(s)
    return " ".join(runs)


def extract_pptx_segments(path: str, source: str) -> list[dict]:
    """One segment per slide: title + tables + SmartArt text, tagged by cadre."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk(shapes, out):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes, out)
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    line = " | ".join(cells).strip(" |")
                    if line:
                        out.append(line)
            elif shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text.strip())

    prs = Presentation(path)
    segments: list[dict] = []
    for slide in prs.slides:
        label_parts: list[str] = []
        walk(slide.shapes, label_parts)  # titles, labels, tables
        smart = _slide_smartart_text(slide)
        # Detect the cadre from the slide's title/labels only — the SmartArt body
        # often mentions other cadres in passing and would blur the tag.
        cadres = _detect_cadres("\n".join(label_parts))
        parts = list(label_parts)
        if smart:
            parts.append(smart)
        text = "\n".join(parts).strip()
        if len(text) < 15:
            continue  # skip near-empty divider/section slides
        segments.append({"source": source, "text": text, "cadres": cadres})
    return segments


def load_segments(directory: str) -> list[dict]:
    """Return segments {source, text, cadres}. PDF/DOCX -> one file-level segment;
    PPTX -> one segment per slide with slide-level cadre tagging."""
    segments: list[dict] = []
    for fname in sorted(os.listdir(directory)):
        fpath = os.path.join(directory, fname)
        if not os.path.isfile(fpath):
            continue
        low = fname.lower()
        if low.endswith(".pdf"):
            print(f"  Reading PDF: {fname}")
            segments.append({"source": fname, "text": extract_pdf(fpath),
                             "cadres": SOURCE_CADRES.get(fname, ["all"])})
        elif low.endswith(".docx"):
            print(f"  Reading DOCX: {fname}")
            segments.append({"source": fname, "text": extract_docx(fpath),
                             "cadres": SOURCE_CADRES.get(fname, ["all"])})
        elif low.endswith(".pptx"):
            print(f"  Reading PPTX: {fname}")
            segs = extract_pptx_segments(fpath, fname)
            print(f"    {len(segs)} slide segment(s) extracted")
            segments.extend(segs)
    return segments


def chunk_segments(segments: list[dict]) -> tuple[list[str], list[dict]]:
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        encoding_name="cl100k_base",
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
    )

    all_texts: list[str] = []
    all_meta: list[dict] = []
    per_source: dict[str, int] = {}

    for seg in segments:
        source, text = seg["source"], seg["text"]
        cadres = seg.get("cadres", ["all"])
        if not text.strip():
            print(f"  WARNING: {source} yielded no extractable text.")
            continue
        for chunk in splitter.split_text(text):
            n = per_source.get(source, 0)
            per_source[source] = n + 1
            all_texts.append(chunk)
            all_meta.append({
                "chunk_id": f"{source}::{n}",
                "text": chunk,
                "source": source,
                "cadres": cadres,
            })

    for source, cnt in per_source.items():
        print(f"  {source}: {cnt} chunks")
    return all_texts, all_meta


def embed_chunks(chunks: list[str]) -> list[list[float]]:
    all_embeddings: list[list[float]] = []
    total_batches = (len(chunks) + EMBED_BATCH_SIZE - 1) // EMBED_BATCH_SIZE

    for i in range(0, len(chunks), EMBED_BATCH_SIZE):
        batch = chunks[i : i + EMBED_BATCH_SIZE]
        resp = client.embeddings.create(model=EMBED_MODEL, input=batch)
        all_embeddings.extend([d.embedding for d in resp.data])
        batch_num = i // EMBED_BATCH_SIZE + 1
        print(f"  Embedded batch {batch_num}/{total_batches} ({len(all_embeddings)}/{len(chunks)} chunks)")

    return all_embeddings


def build_and_save_index(embeddings: list[list[float]], metadata: list[dict]) -> None:
    emb_array = np.array(embeddings, dtype=np.float32)
    faiss.normalize_L2(emb_array)

    dim = emb_array.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(emb_array)

    os.makedirs(FAISS_INDEX_PATH, exist_ok=True)
    faiss.write_index(index, os.path.join(FAISS_INDEX_PATH, "index.faiss"))
    with open(os.path.join(FAISS_INDEX_PATH, "metadata.json"), "w", encoding="utf-8") as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)

    print(f"\nFAISS index saved -> {FAISS_INDEX_PATH}/")
    print(f"  Vectors: {index.ntotal}  Dimensions: {dim}")


def main() -> None:
    print(f"Loading documents from: {TRAINING_MATERIALS_DIR}")
    segments = load_segments(TRAINING_MATERIALS_DIR)
    if not segments:
        print("No PDF/DOCX/PPTX files found. Copy the source files to data/training_materials/ and re-run.")
        sys.exit(1)
    sources = {s["source"] for s in segments}
    print(f"Loaded {len(sources)} document(s), {len(segments)} segment(s)\n")

    print("Chunking documents...")
    texts, metadata = chunk_segments(segments)
    if not texts:
        print("ERROR: All documents yielded zero text chunks.")
        print("The PDFs are likely scanned/image-based. Install pymupdf and use the fallback in CLAUDE.md.")
        sys.exit(1)
    print(f"Total chunks: {len(texts)}\n")

    print("Embedding chunks (may take several minutes for large documents)...")
    embeddings = embed_chunks(texts)
    print()

    print("Building FAISS index...")
    build_and_save_index(embeddings, metadata)
    print(f"\nDone. {len(sources)} documents -> {len(texts)} chunks indexed.")


if __name__ == "__main__":
    main()
